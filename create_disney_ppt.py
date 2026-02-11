#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建东京迪士尼家庭游玩攻略PPT
基于大纲文件：东京迪士尼家庭游玩攻略PPT大纲.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

def create_disney_ppt():
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片宽高比（16:9）
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 幻灯片1：封面页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "东京迪士尼家庭游玩全攻略"
    subtitle.text = "2026年2月 · 家庭欢乐之旅\n\n制作人：Helen\n日期：2026年2月"
    
    # 设置标题样式
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)  # 迪士尼蓝
    
    # 幻灯片2：目录
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容
    title2 = slide2.shapes.title
    title2.text = "目录"
    title2.text_frame.paragraphs[0].font.size = Pt(36)
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 51)  # 迪士尼红
    
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.clear()  # 清空默认文本
    
    # 添加目录项
    items = [
        "1. 行前准备篇",
        "2. 交通指南篇", 
        "3. 门票攻略篇",
        "4. 迪士尼乐园攻略",
        "5. 迪士尼海洋攻略",
        "6. 必玩项目推荐",
        "7. 餐饮美食指南",
        "8. 游行演出时刻",
        "9. 购物纪念品攻略",
        "10. 实用贴士汇总"
    ]
    
    for item in items:
        p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.bold = True
        p.space_after = Pt(12)
    
    # 幻灯片3：行前准备篇
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "行前准备篇"
    title3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 0)  # 绿色
    
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.clear()
    
    p3 = tf3.add_paragraph()
    p3.text = "重要事项清单"
    p3.font.size = Pt(28)
    p3.font.bold = True
    
    items3 = [
        "签证：确认日本旅游签证有效期",
        "机票：提前预订，建议选择成田或羽田机场",
        "住宿：迪士尼官方酒店（提前3-6个月预订）",
        "天气：2月东京气温5-10°C，准备保暖衣物",
        "必备物品：护照、身份证件复印件、充电宝、转换插头等"
    ]
    
    for item in items3:
        p = tf3.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.level = 0
    
    # 幻灯片4：交通指南篇
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "交通指南篇"
    
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.clear()
    
    p4 = tf4.add_paragraph()
    p4.text = "从机场到迪士尼"
    p4.font.size = Pt(24)
    p4.font.bold = True
    
    items4_1 = [
        "成田机场→迪士尼：利木津巴士（直达，约90分钟）",
        "成田机场→迪士尼：JR成田特快→东京站→换乘JR京叶线",
        "羽田机场→迪士尼：京急线→品川站→换乘JR京叶线",
        "羽田机场→迪士尼：利木津巴士（约60分钟）"
    ]
    
    for item in items4_1:
        p = tf4.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.level = 0
    
    p4_2 = tf4.add_paragraph()
    p4_2.text = "\n园区间交通"
    p4_2.font.size = Pt(24)
    p4_2.font.bold = True
    
    items4_2 = [
        "迪士尼度假区线（单轨电车）：连接乐园、海洋、酒店、购物区",
        "一日券：成人660日元，儿童330日元"
    ]
    
    for item in items4_2:
        p = tf4.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.level = 0
    
    # 幻灯片5：门票攻略篇
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "门票攻略篇"
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.clear()
    
    items5 = [
        "一日护照：单园游玩",
        "两日护照：可分别游玩两园",
        "三日/四日魔法护照：多日游玩更优惠",
        "星光护照：下午3点后入园（价格优惠）",
        "",
        "购买渠道：",
        "• 官方渠道：东京迪士尼官网、官方APP",
        "• 其他渠道：日本便利店（Lawson等）、旅行社代订",
        "",
        "重要提醒：",
        "• 提前1-2个月购买，旺季票源紧张",
        "• 儿童票：4-11岁，4岁以下免费",
        "• 打印门票或保存电子票二维码"
    ]
    
    for item in items5:
        if not item.strip():  # 跳过空行
            p = tf5.add_paragraph()
            p.text = ""
            continue
        p = tf5.add_paragraph()
        p.text = item
        if "：" in item and "•" not in item:
            p.font.size = Pt(22)
            p.font.bold = True
        elif "•" in item:
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.font.size = Pt(20)
    
    # 幻灯片6：迪士尼乐园攻略
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "迪士尼乐园攻略"
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.clear()
    
    items6 = [
        "园区概况：",
        "• 主题：经典迪士尼童话世界",
        "• 适合人群：亲子家庭、迪士尼粉丝",
        "• 开放时间：通常8:00-22:00（具体以官网为准）",
        "",
        "分区介绍：",
        "1. 世界市集：购物、餐饮主街",
        "2. 明日乐园：太空山、星际旅行",
        "3. 卡通城：适合低龄儿童",
        "4. 梦幻乐园：灰姑娘城堡、小小世界",
        "5. 动物天地：飞溅山",
        "6. 西部乐园：巨雷山",
        "7. 探险乐园：加勒比海盗"
    ]
    
    for item in items6:
        if not item.strip():  # 跳过空行
            p = tf6.add_paragraph()
            p.text = ""
            continue
        p = tf6.add_paragraph()
        p.text = item
        if "：" in item and not item[0].isdigit():
            p.font.size = Pt(22)
            p.font.bold = True
        elif item[0].isdigit():
            p.font.size = Pt(18)
            p.level = 0
        elif "•" in item:
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.font.size = Pt(20)
    
    # 幻灯片7：迪士尼海洋攻略
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "迪士尼海洋攻略"
    
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.clear()
    
    items7 = [
        "园区概况：",
        "• 全球独有：唯一以海洋为主题的迪士尼乐园",
        "• 适合人群：青少年、成年人、寻求刺激者",
        "• 开放时间：通常8:00-22:00",
        "",
        "分区介绍：",
        "1. 地中海港湾：入口区域，仿意大利渔村",
        "2. 美国海滨：惊魂古塔、玩具总动员疯狂游戏屋",
        "3. 发现港：海底巡游艇、风暴骑士",
        "4. 失落河三角洲：印第安纳琼斯冒险",
        "5. 阿拉伯海岸：辛巴达传奇之旅",
        "6. 美人鱼礁湖：适合低龄儿童",
        "7. 神秘岛：地心探险之旅、海底两万里"
    ]
    
    for item in items7:
        if not item.strip():  # 跳过空行
            p = tf7.add_paragraph()
            p.text = ""
            continue
        p = tf7.add_paragraph()
        p.text = item
        if "：" in item and not item[0].isdigit():
            p.font.size = Pt(22)
            p.font.bold = True
        elif item[0].isdigit():
            p.font.size = Pt(18)
            p.level = 0
        elif "•" in item:
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.font.size = Pt(20)
    
    # 幻灯片8：必玩项目推荐
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "必玩项目推荐"
    
    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.clear()
    
    items8 = [
        "迪士尼乐园Top 5：",
        "1. 巨雷山（西部乐园）：家庭过山车，刺激度适中",
        "2. 飞溅山（动物天地）：水上漂流，有落差惊喜",
        "3. 太空山（明日乐园）：室内过山车，星空体验",
        "4. 美女与野兽城堡（梦幻乐园）：新园区，浪漫之旅",
        "5. 小小世界（梦幻乐园）：经典游船，儿童最爱",
        "",
        "迪士尼海洋Top 5：",
        "1. 玩具总动员疯狂游戏屋（美国海滨）：互动射击游戏",
        "2. 惊魂古塔（美国海滨）：跳楼机，刺激度高",
        "3. 印第安纳琼斯冒险（失落河三角洲）：沉浸式冒险",
        "4. 地心探险之旅（神秘岛）：火山过山车",
        "5. 海底两万里（神秘岛）：潜水艇探险"
    ]
    
    for item in items8:
        if not item.strip():  # 跳过空行
            p = tf8.add_paragraph()
            p.text = ""
            continue
        p = tf8.add_paragraph()
        p.text = item
        if "：" in item and not item[0].isdigit():
            p.font.size = Pt(22)
            p.font.bold = True
        elif item[0].isdigit():
            p.font.size = Pt(18)
            p.level = 0
        else:
            p.font.size = Pt(20)
    
    # 幻灯片9：餐饮美食指南
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "餐饮美食指南"
    
    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.clear()
    
    items9 = [
        "餐厅类型：",
        "• 自助餐厅：品类丰富，适合家庭",
        "• 主题餐厅：角色互动，体验独特",
        "• 快餐小吃：便捷，节省时间",
        "• 特色小吃：米奇形状食品、限定美食",
        "",
        "推荐餐厅：",
        "• 迪士尼乐园：皇冠餐厅、红心女王宴会大厅、明日乐园舞台餐厅",
        "• 迪士尼海洋：苏丹的绿洲、麦哲伦餐厅、船长的厨房",
        "",
        "预算参考：",
        "• 成人套餐：1500-3000日元",
        "• 儿童套餐：800-1500日元",
        "• 小吃：500-1000日元"
    ]
    
    for item in items9:
        if not item.strip():  # 跳过空行
            p = tf9.add_paragraph()
            p.text = ""
            continue
        p = tf9.add_paragraph()
        p.text = item
        if "：" in item and "•" not in item:
            p.font.size = Pt(22)
            p.font.bold = True
        elif "•" in item:
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.font.size = Pt(20)
    
    # 幻灯片10：游行演出时刻
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "游行演出时刻"
    
    content10 = slide10.placeholders[1]
    tf10 = content10.text_frame
    tf10.clear()
    
    items10 = [
        "日间游行：",
        "• 迪士尼乐园：'幸福在这里'大游行（通常14:00）",
        "• 迪士尼海洋：'水上迎宾'表演（地中海港湾）",
        "",
        "夜间娱乐：",
        "• 迪士尼乐园：'东京迪士尼乐园电子大游行～梦之光'、烟花秀",
        "• 迪士尼海洋：'Fantasmic!'水上秀、'彩色世界'灯光秀",
        "",
        "观看贴士：",
        "• 提前1-2小时占位（热门区域）",
        "• 购买预留座位票（如有）",
        "• 下载官方APP查看当日时刻表"
    ]
    
    for item in items10:
        if not item.strip():  # 跳过空行
            p = tf10.add_paragraph()
            p.text = ""
            continue
        p = tf10.add_paragraph()
        p.text = item
        if "：" in item and "•" not in item:
            p.font.size = Pt(22)
            p.font.bold = True
        elif "•" in item:
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.font.size = Pt(20)
    
    # 幻灯片11：购物纪念品攻略
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    title11.text = "购物纪念品攻略"
    
    content11 = slide11.placeholders[1]
    tf11 = content11.text_frame
    tf11.clear()
    
    items11 = [
        "必买商品：",
        "• 迪士尼限定：东京迪士尼独家商品",
        "• 角色周边：米奇、达菲、星黛露等",
        "• 季节性商品：2月可能有情人节限定",
        "• 实用商品：服装、文具、家居用品",
        "",
        "购物地点：",
        "• 世界市集（迪士尼乐园）：最大购物区",
        "• 迪士尼海洋：每个港口都有特色商店",
        "• 伊克斯皮儿莉购物中心：园区外大型商场",
        "• 酒店商店：部分商品酒店专售",
        "",
        "退税信息：",
        "• 购物满5000日元可办理退税",
        "• 需出示护照原件",
        "• 部分商店直接免税"
    ]
    
    for item in items11:
        if not item.strip():  # 跳过空行
            p = tf11.add_paragraph()
            p.text = ""
            continue
        p = tf11.add_paragraph()
        p.text = item
        if "：" in item and "•" not in item:
            p.font.size = Pt(22)
            p.font.bold = True
        elif "•" in item:
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.font.size = Pt(20)
    
    # 幻灯片12：实用贴士汇总
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    title12.text = "实用贴士汇总"
    
    content12 = slide12.placeholders[1]
    tf12 = content12.text_frame
    tf12.clear()
    
    items12 = [
        "节省时间技巧：",
        "1. 早到原则：开园前30分钟到达",
        "2. FastPass/尊享卡：合理利用快速通道",
        "3. 单人通道：部分项目可走单人通道",
        "4. 错峰游玩：午餐时间、游行时间排队较短",
        "",
        "家庭便利设施：",
        "• 婴儿车租赁：1000日元/天",
        "• 母婴室：各园区均有",
        "• 储物柜：小型300日元，大型500-700日元",
        "• 轮椅租赁：免费（需押金）",
        "",
        "其他贴士：",
        "• 下载官方APP：查看排队时间、地图",
        "• 穿着舒适鞋子：日行2万步以上",
        "• 带好充电宝：拍照、APP使用耗电快",
        "• 学习简单日语：问路、点餐更便利"
    ]
    
    for item in items12:
        if not item.strip():  # 跳过空行
            p = tf12.add_paragraph()
            p.text = ""
            continue
        p = tf12.add_paragraph()
        p.text = item
        if "：" in item and not item[0].isdigit():
            p.font.size = Pt(22)
            p.font.bold = True
        elif item[0].isdigit():
            p.font.size = Pt(18)
            p.level = 0
        elif "•" in item:
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.font.size = Pt(20)
    
    # 幻灯片13：预算估算参考
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    title13 = slide13.shapes.title
    title13.text = "预算估算参考"
    
    content13 = slide13.placeholders[1]
    tf13 = content13.text_frame
    tf13.clear()
    
    items13 = [
        "三日两晚家庭游（2大1小）：",
        "• 机票：3000-5000元/人",
        "• 住宿：迪士尼官方酒店2000-4000元/晚",
        "• 门票：两日护照约1500元/人",
        "• 餐饮：1500-2500元/天",
        "• 购物：视个人需求",
        "• 交通：500-1000元",
        "• 总计：15000-25000元（仅供参考）"
    ]
    
    for item in items13:
        if not item.strip():  # 跳过空行
            p = tf13.add_paragraph()
            p.text = ""
            continue
        p = tf13.add_paragraph()
        p.text = item
        if "：" in item and "•" not in item:
            p.font.size = Pt(22)
            p.font.bold = True
        elif "•" in item:
            p.font.size = Pt(20)
            p.level = 0
        else:
            p.font.size = Pt(20)
    
    # 幻灯片14：结束页
    slide14 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片
    title14 = slide14.shapes.title
    subtitle14 = slide14.placeholders[1]
    
    title14.text = "开启魔法之旅！"
    title14.text_frame.paragraphs[0].font.size = Pt(48)
    title14.text_frame.paragraphs[0].font.bold = True
    title14.text_frame.paragraphs[0].font.color.rgb = RGBColor(153, 0, 153)  # 紫色
    
    subtitle14.text = "祝您和家人在东京迪士尼度过难忘的欢乐时光！\n\n制作人：Helen\n联系方式：helen@example.com\n\n感谢观看！"
    
    # 保存PPT文件
    output_path = "/Users/yangyan/.openclaw/workspace/东京迪士尼家庭游玩攻略.pptx"
    prs.save(output_path)
    
    print(f"PPT已创建：{output_path}")
    print(f"总幻灯片数：{len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    try:
        ppt_path = create_disney_ppt()
        print("PPT创建成功！")
    except Exception as e:
        print(f"创建PPT时出错：{e}")
        import traceback
        traceback.print_exc()